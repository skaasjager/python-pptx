# encoding: utf-8

"""
Gherkin step implementations for shape-related features.
"""

import os

from behave import given, when, then
from hamcrest import assert_that, equal_to, is_

from pptx import Presentation
from pptx.constants import MSO_AUTO_SHAPE_TYPE as MAST, MSO
from pptx.util import Inches


def absjoin(*paths):
    return os.path.abspath(os.path.join(*paths))

thisdir = os.path.split(__file__)[0]
scratch_dir = absjoin(thisdir, '../_scratch')
saved_pptx_path = absjoin(scratch_dir, 'test_out.pptx')

test_text = "python-pptx was here!"


# given ===================================================

@given('I have a reference to a chevron shape')
def step_given_ref_to_chevron_shape(context):
    context.prs = Presentation()
    blank_slidelayout = context.prs.slidelayouts[6]
    shapes = context.prs.slides.add_slide(blank_slidelayout).shapes
    x = y = cx = cy = 914400
    context.chevron_shape = shapes.add_shape(MAST.CHEVRON, x, y, cx, cy)


# when ====================================================

@when("I add a text box to the slide's shape collection")
def step_when_add_text_box(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(1.00))
    sp = shapes.add_textbox(x, y, cx, cy)
    sp.text = test_text


@when("I add an auto shape to the slide's shape collection")
def step_when_add_auto_shape(context):
    shapes = context.sld.shapes
    x, y = (Inches(1.00), Inches(2.00))
    cx, cy = (Inches(3.00), Inches(4.00))
    sp = shapes.add_shape(MAST.ROUNDED_RECTANGLE, x, y, cx, cy)
    sp.text = test_text


@when("I set the first adjustment value to 0.15")
def step_when_set_first_adjustment_value(context):
    context.chevron_shape.adjustments[0] = 0.15


# then ====================================================

@then('the auto shape appears in the slide')
def step_then_auto_shape_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    sp = prs.slides[0].shapes[0]
    sp_text = sp.textframe.paragraphs[0].runs[0].text
    assert_that(sp.shape_type, is_(equal_to(MSO.AUTO_SHAPE)))
    assert_that(sp.auto_shape_type, is_(equal_to(MAST.ROUNDED_RECTANGLE)))
    assert_that(sp_text, is_(equal_to(test_text)))


@then('the chevron shape appears with a less acute arrow head')
def step_then_chevron_shape_appears_with_less_acute_arrow_head(context):
    chevron = Presentation(saved_pptx_path).slides[0].shapes[0]
    assert_that(chevron.adjustments[0], is_(equal_to(0.15)))


@then('the text box appears in the slide')
def step_then_text_box_appears_in_slide(context):
    prs = Presentation(saved_pptx_path)
    textbox = prs.slides[0].shapes[0]
    textbox_text = textbox.textframe.paragraphs[0].runs[0].text
    assert_that(textbox_text, is_(equal_to(test_text)))